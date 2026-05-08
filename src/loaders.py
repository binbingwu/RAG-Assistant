from __future__ import annotations

import json
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

from bs4 import BeautifulSoup
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader


SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".pptx",
    ".xlsx",
    ".txt",
    ".md",
    ".csv",
    ".json",
    ".js",
    ".html",
    ".htm",
    ".png",
}


@dataclass
class LoadedDocument:
    path: str
    file_name: str
    extension: str
    text: str


def iter_supported_files(source_dir: Path) -> Iterable[Path]:
    for path in source_dir.rglob("*"):
        if path.is_file() and path.suffix.lower() in SUPPORTED_EXTENSIONS:
            yield path


def load_document(path: Path) -> LoadedDocument:
    ext = path.suffix.lower()
    if ext == ".pdf":
        text = _load_pdf(path)
    elif ext == ".docx":
        text = _load_docx(path)
    elif ext == ".pptx":
        text = _load_pptx(path)
    elif ext == ".xlsx":
        text = _load_xlsx(path)
    elif ext in {".html", ".htm"}:
        text = _load_html(path)
    elif ext == ".json":
        text = _load_json(path)
    elif ext == ".png":
        text = ""
    else:
        text = _read_text(path)

    return LoadedDocument(
        path=str(path),
        file_name=path.name,
        extension=ext,
        text=_normalize_text(text),
    )


def _load_pdf(path: Path) -> str:
    reader = PdfReader(str(path))
    pages = []
    for idx, page in enumerate(reader.pages, start=1):
        page_text = page.extract_text() or ""
        pages.append(f"\n[Page {idx}]\n{page_text}")
    return "\n".join(pages)


def _load_docx(path: Path) -> str:
    doc = Document(str(path))
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _load_pptx(path: Path) -> str:
    prs = Presentation(str(path))
    parts = []
    for idx, slide in enumerate(prs.slides, start=1):
        slide_parts = [f"[Slide {idx}]"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_parts.append(shape.text)
        parts.append("\n".join(slide_parts))
    return "\n\n".join(parts)


def _load_xlsx(path: Path) -> str:
    wb = load_workbook(filename=str(path), read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(f"[Sheet: {ws.title}]")
        for row in ws.iter_rows(values_only=True):
            values = [str(value).strip() for value in row if value is not None]
            if values:
                parts.append(" | ".join(values))
    return "\n".join(parts)


def _load_html(path: Path) -> str:
    soup = BeautifulSoup(_read_text(path), "html.parser")
    return soup.get_text("\n")


def _load_json(path: Path) -> str:
    try:
        data = json.loads(_read_text(path))
        return json.dumps(data, ensure_ascii=False, indent=2)
    except json.JSONDecodeError:
        return _read_text(path)


def _read_text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def _normalize_text(text: str) -> str:
    lines = [line.strip() for line in text.replace("\r", "\n").split("\n")]
    return "\n".join(line for line in lines if line)
